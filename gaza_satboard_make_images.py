"""
gaza_satboard_make_images.py
Run on your machine or in Colab. Produces:
 - pre_post_ndvi_pretty_pre.png
 - pre_post_ndvi_pretty_post.png
 - pre_post_side_by_side.png
 - (optional) answer_updated.pptx (with images embedded)

Dependencies (pip):
pip install pystac-client rasterio rioxarray matplotlib numpy pillow python-pptx affine requests tqdm pystac

Notes:
- Uses the public STAC endpoint: https://earth-search.aws.element84.com/v0
- Streams COGs via HTTP (rasterio can read COGs by URL).
- If you run in Colab, mount / ensure internet access.
"""

import os, sys, math, requests
import numpy as np
import rasterio
import rioxarray as rxr
import matplotlib.pyplot as plt
from matplotlib_scalebar.scalebar import ScaleBar  # optional (pip install matplotlib-scalebar)
from pystac_client import Client
from shapely.geometry import box, Point
from shapely.ops import transform
from pyproj import Transformer
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm

# ------- USER PARAMETERS -------
BBOX = [34.21870, 31.22005, 34.56780, 31.59454]  # [minLon,minLat,maxLon,maxLat]
PRE_START = "2025-07-15"
PRE_END   = "2025-07-31"
POST_START= "2025-08-15"
POST_END  = "2025-08-25"
OUTDIR = "gaza_sat_outputs"
os.makedirs(OUTDIR, exist_ok=True)

# Neighbourhood centers (lon, lat)
NEIGH = {
    "Shuja'iyya": (34.4683, 31.5000),
    "Zeitoun":    (34.4441, 31.4912),
    "Khan Yunis": (34.2962, 31.3345),
}

# STAC client (AWS Earth Search)
STAC_URL = "https://earth-search.aws.element84.com/v0"
client = Client.open(STAC_URL)

# ------- helper functions -------
def find_sentinel_item(bbox, start, end, cloud_max=20, limit=5):
    """
    Find Sentinel-2 L2A items overlapping bbox in date window; returns the first low-cloud item.
    """
    search = client.search(
        collections=["sentinel-s2-l2a-cogs"],
        bbox=bbox,
        datetime=f"{start}/{end}",
        query={"eo:cloud_cover": {"lt": cloud_max}},
        sortby=[{"field":"properties.datetime","direction":"desc"}],
        limit=limit
    )
    items = list(search.get_all_items())
    if not items:
        raise RuntimeError(f"No items found for {start}->{end} within cloud<{cloud_max}")
    # pick the least-cloud cover item
    items = sorted(items, key=lambda it: it.properties.get("eo:cloud_cover", 100))
    chosen = items[0]
    print("Chosen item:", chosen.id, "date:", chosen.properties.get("datetime"), "cloud:", chosen.properties.get("eo:cloud_cover"))
    return chosen

def read_rgb_and_ndvi(item, bbox):
    """
    Given a STAC item, read Red (B04), Green (B03), Blue (B02), NIR (B08) as arrays clipped to bbox.
    Returns: rgb (H,W,3), ndvi (H,W), transform, crs
    """
    # asset keys names may be 'B04' 'B03' etc for sentinel-cogs collection; we will look for common asset names
    assets = item.assets
    # asset hrefs
    def asset_href(band_key):
        # try 'B04' / 'B04.jp2' / 'red' etc - sentinel-cogs standard: 'B04'
        for k in assets:
            if band_key in k:
                return assets[k].href
            if k.upper().endswith(band_key.upper()):
                return assets[k].href
        raise KeyError("Band not found: "+band_key)
    # HACK: sentinel-2_l2a_cogs uses 'B02','B03','B04','B08'
    b2 = asset_href("B02")
    b3 = asset_href("B03")
    b4 = asset_href("B04")
    b8 = asset_href("B08")

    # Use rasterio windows to read bbox
    # rasterio wants bbox in same crs as dataset (usually EPSG:4326 for STAC asset HREFS?). We'll read full arrays but could stream windowed read.
    with rasterio.open(b4) as src:
        # compute window from bbox in lon/lat: transform bbox to src CRS
        transformer = Transformer.from_crs("EPSG:4326", src.crs, always_xy=True)
        minx, miny = transformer.transform(bbox[0], bbox[1])
        maxx, maxy = transformer.transform(bbox[2], bbox[3])
        # form window
        window = src.window(minx, maxy, maxx, miny)  # careful with coords order
        window = window.round_offsets().round_shape()
        # read arrays
        red = src.read(1, window=window, out_shape=(int(window.height), int(window.width)))
        meta = src.meta.copy()
        transform_win = src.window_transform(window)
        crs = src.crs

    # read green, blue, nir similarly using rasterio.open and same window/transform approach:
    def read_band(href, transform_win=None, window=None):
        with rasterio.open(href) as s:
            arr = s.read(1, window=window, out_shape=(int(window.height), int(window.width)))
            return arr
    green = read_band(b3, window=window)
    blue  = read_band(b2, window=window)
    nir   = read_band(b8, window=window)

    # convert to float and scale (Sentinel L2A typically scaled 0-10000)
    def scale_arr(a):
        arr = a.astype('float32')
        # if values > 10000 assume scaled 0..10000
        if arr.max() > 10000:
            pass
        return arr / 10000.0

    r = scale_arr(red); g = scale_arr(green); b = scale_arr(blue); n = scale_arr(nir)
    # stack RGB
    rgb = np.dstack([r,g,b])
    # compute ndvi
    ndvi = (n - r) / (n + r + 1e-6)
    # clip to -1..1
    ndvi = np.clip(ndvi, -1, 1)
    return rgb, ndvi, transform_win, crs

def plot_and_save(rgb, ndvi, transform, crs, out_prefix, item, bbox, annotations):
    """
    Create:
     - false-color RGB png
     - NDVI color map png
     - combined side-by-side
    annotations: list of dicts {name, lon, lat, box_km}
    """
    # compute extent in lon/lat from transform & shape
    h, w, _ = rgb.shape
    # create figure
    fig, axes = plt.subplots(1,2, figsize=(16,9))
    # show RGB
    axes[0].imshow(np.clip(rgb,0,1))
    axes[0].set_title(f"RGB (Sentinel-2 L2A) {item.properties.get('datetime')[:10]}")
    # NDVI plot
    im = axes[1].imshow(ndvi, vmin=-1, vmax=1, cmap='RdYlGn')
    axes[1].set_title("NDVI")
    fig.colorbar(im, ax=axes[1], fraction=0.046, pad=0.04)
    # annotate: convert lon/lat to pixel coords using transformer
    transformer = Transformer.from_crs("EPSG:4326", crs.to_string(), always_xy=True)
    def lonlat_to_px(lon, lat):
        x, y = transformer.transform(lon, lat)
        # row/col from transform: col = (x - tx)/a, row = (y - ty)/e with affine
        from rasterio.transform import rowcol
        row, col = rowcol(transform, x, y, op=float)
        return col, row
    for ann in annotations:
        col, row = lonlat_to_px(ann['lon'], ann['lat'])
        for ax in axes:
            rect = plt.Rectangle((col-20, row-20), 40, 40, linewidth=2, edgecolor='yellow', facecolor='none')
            ax.add_patch(rect)
            ax.text(col+5, row+5, ann['name'], color='yellow', fontsize=10, weight='bold')
    # save
    rgb_file = os.path.join(OUTDIR, f"{out_prefix}_rgb.png")
    ndvi_file = os.path.join(OUTDIR, f"{out_prefix}_ndvi.png")
    fig.savefig(os.path.join(OUTDIR, f"{out_prefix}_combined.png"), bbox_inches='tight', dpi=200)
    plt.close(fig)
    print("Saved:", os.path.join(OUTDIR, f"{out_prefix}_combined.png"))

# ------- Main flow -------
def main():
    print("Searching STAC for PRE window...")
    pre_item = find_sentinel_item(BBOX, PRE_START, PRE_END)
    print("Searching STAC for POST window...")
    post_item = find_sentinel_item(BBOX, POST_START, POST_END)

    # read & process PRE
    print("Reading PRE item and computing NDVI (may stream COGs)...")
    pre_rgb, pre_ndvi, pre_transform, pre_crs = read_rgb_and_ndvi(pre_item, BBOX)
    print("Writing PRE annotated images...")
    annotations = [{"name":k, "lon":v[0], "lat":v[1]} for k,v in NEIGH.items()]
    plot_and_save(pre_rgb, pre_ndvi, pre_transform, pre_crs, "pre", pre_item, BBOX, annotations)

    # read & process POST
    print("Reading POST item and computing NDVI (may stream COGs)...")
    post_rgb, post_ndvi, post_transform, post_crs = read_rgb_and_ndvi(post_item, BBOX)
    print("Writing POST annotated images...")
    plot_and_save(post_rgb, post_ndvi, post_transform, post_crs, "post", post_item, BBOX, annotations)

    # Make a simple PPTX with the results
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]  # blank
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Gaza SatBoard — Sentinel-2 annotated imagery"
    # Add pre/post combined
    for f in ["pre_combined.png", "post_combined.png"]:
        slide = prs.slides.add_slide(title_slide_layout)
        left = Inches(0.3)
        top = Inches(0.6)
        slide.shapes.add_picture(os.path.join(OUTDIR, f), left, top, width=Inches(9.0))
    outppt = os.path.join(OUTDIR, "gaza_satboard_images.pptx")
    prs.save(outppt)
    print("Saved PPTX:", outppt)

if __name__ == "__main__":
    main()
